"""Build src/data/courseContent.json from documentation/."""
import json
import re
from pathlib import Path

from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
DOC_DIR = ROOT / "documentation"
OUT = ROOT / "src" / "data" / "courseContent.json"


def slugify(name: str) -> str:
    s = re.sub(r"\.[^.]+$", "", name)
    s = re.sub(r"[^a-zA-Z0-9]+", "-", s).strip("-").lower()
    return s or "doc"


def extract_docx_paragraphs(path: Path) -> list[str]:
    doc = Document(path)
    out: list[str] = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            out.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                out.append(" | ".join(cells))
    return out


def extract_pptx_slides(path: Path) -> list[dict]:
    prs = Presentation(path)
    slides_out: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if t:
                chunks.append(t)
        full = "\n".join(chunks)
        lines = [ln.strip() for ln in full.splitlines() if ln.strip()]
        slides_out.append({"number": i, "lines": lines, "raw": full})
    return slides_out


def doc_title(filename: str) -> str:
    base = re.sub(r"\.[^.]+$", "", filename)
    mapping = {
        "Appendix F HW Questions": "Appendix F: Homework Questions",
        "Getting Started 012720261": "Getting Started (Course Intro)",
        "Internet Basics_03042026": "Internet Basics & Networking",
        "Midterm topics & format Spring 2026v1": "Midterm: Topics & Exam Format",
        "CH01-CompSec5e_02022026": "Chapter 1: Security Overview",
        "CH02-CompSec5e_0304026": "Chapter 2: Cryptographic Tools",
        "CH03-CompSec5e_03032026a": "Chapter 3: User Authentication",
        "CH06-CompSec5e_04072026": "Chapter 6: Malicious Software",
        "CH07-CompSec5e_03242026": "Chapter 7: Denial-of-Service Attacks",
        "CH08-CompSec5e_03312026": "Chapter 8: Intrusion Detection",
        "CH09-CompSec5e_04022026": "Chapter 9: Firewalls & IPS",
        "CH10-CompSec3e040824": "Chapter 10: Software Security (3rd Ed. slides)",
        "CH10-CompSec5e_accessible": "Chapter 10: Software Security (5th Ed., accessible)",
    }
    return mapping.get(base, base.replace("_", " "))


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    items: list[dict] = []

    for path in sorted(DOC_DIR.iterdir(), key=lambda x: x.name.lower()):
        suf = path.suffix.lower()
        if suf == ".docx":
            paras = extract_docx_paragraphs(path)
            items.append(
                {
                    "id": slugify(path.name),
                    "title": doc_title(path.name),
                    "sourceFile": path.name,
                    "kind": "document",
                    "paragraphs": paras,
                }
            )
        elif suf == ".pptx":
            slides = extract_pptx_slides(path)
            items.append(
                {
                    "id": slugify(path.name),
                    "title": doc_title(path.name),
                    "sourceFile": path.name,
                    "kind": "slides",
                    "slides": slides,
                }
            )

    payload = {
        "version": 1,
        "generatedNote": "Auto-built from documentation/; do not hand-edit for content—re-run scripts/build_course_json.py",
        "items": items,
    }
    OUT.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Wrote {OUT} with {len(items)} items")


if __name__ == "__main__":
    main()
