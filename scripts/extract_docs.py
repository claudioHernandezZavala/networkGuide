"""One-off extraction of documentation text for study guide."""
import os
from pathlib import Path

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

ROOT = Path(__file__).resolve().parents[1]
DOC_DIR = ROOT / "documentation"
OUT = ROOT / "scripts" / "extracted_raw.txt"


def extract_docx(path: Path) -> str:
    if not Document:
        return ""
    doc = Document(path)
    lines = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            lines.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def extract_pptx(path: Path) -> str:
    if not Presentation:
        return ""
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---\n")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if t:
                lines.append(t)
    return "\n".join(lines)


def main():
    parts = []
    for p in sorted(DOC_DIR.iterdir(), key=lambda x: x.name.lower()):
        if p.suffix.lower() == ".docx":
            parts.append(f"\n\n######## {p.name} (DOCX)\n")
            parts.append(extract_docx(p))
        elif p.suffix.lower() == ".pptx":
            parts.append(f"\n\n######## {p.name} (PPTX)\n")
            parts.append(extract_pptx(p))
    OUT.write_text("\n".join(parts), encoding="utf-8")
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
